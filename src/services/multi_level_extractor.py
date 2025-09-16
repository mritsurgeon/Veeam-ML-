"""
Multi-Level Data Extraction Service

This service implements the three-level data extraction strategy:
1. File System Level (Metadata Only)
2. File Content Level (Documents/General Files)  
3. Database Level (Structured Data Extraction)
"""

import os
import logging
import json
import sqlite3
import pandas as pd
from typing import Dict, List, Any, Optional, Union
from datetime import datetime
import hashlib
import mimetypes
from pathlib import Path

# Document parsing libraries
try:
    import PyPDF2
    from docx import Document
    import openpyxl
    from pptx import Presentation
    HAS_DOCUMENT_PARSERS = True
except ImportError:
    HAS_DOCUMENT_PARSERS = False
    logging.warning("Document parsing libraries not available. Install: pip install PyPDF2 python-docx openpyxl python-pptx")

# Database libraries
try:
    import sqlite3
    import pymysql
    import psycopg2
    HAS_DB_LIBRARIES = True
except ImportError:
    HAS_DB_LIBRARIES = False
    logging.warning("Database libraries not available. Install: pip install pymysql psycopg2-binary")

logger = logging.getLogger(__name__)

class MultiLevelExtractor:
    """
    Multi-level data extraction service for ML analysis.
    """
    
    def __init__(self, chunk_size: int = 2000, max_file_size: int = 50 * 1024 * 1024):
        """
        Initialize the multi-level extractor.
        
        Args:
            chunk_size: Size of text chunks for processing
            max_file_size: Maximum file size to process (50MB default)
        """
        self.chunk_size = chunk_size
        self.max_file_size = max_file_size
        self.extraction_stats = {
            'files_processed': 0,
            'chunks_created': 0,
            'databases_extracted': 0,
            'errors': 0
        }
    
    def extract_data(self, file_info: Dict[str, Any], unc_path: str, 
                    extraction_level: str = 'auto') -> Dict[str, Any]:
        """
        Extract data from a file based on the specified level.
        
        Args:
            file_info: File information from metadata extraction
            unc_path: UNC path to the file
            extraction_level: Level of extraction ('metadata', 'content', 'database', 'auto')
            
        Returns:
            Dictionary containing extracted data and metadata
        """
        try:
            file_type = file_info.get('file_type', 'unknown')
            file_size = file_info.get('size', 0)
            
            # Auto-determine extraction level if not specified
            if extraction_level == 'auto':
                extraction_level = self._determine_extraction_level(file_type, file_size)
            
            result = {
                'file_info': file_info,
                'extraction_level': extraction_level,
                'extraction_timestamp': datetime.utcnow().isoformat(),
                'extracted_data': None,
                'chunks': [],
                'metadata': {},
                'error': None
            }
            
            if extraction_level == 'metadata':
                result['extracted_data'] = self._extract_metadata_only(file_info)
                
            elif extraction_level == 'content':
                if file_size > self.max_file_size:
                    result['error'] = f"File too large for content extraction: {file_size} bytes"
                else:
                    content_data = self._extract_file_content(file_info, unc_path)
                    result['extracted_data'] = content_data
                    result['chunks'] = self._chunk_text(content_data.get('text', ''))
                    
            elif extraction_level == 'database':
                db_data = self._extract_database_content(file_info, unc_path)
                result['extracted_data'] = db_data
                
            self.extraction_stats['files_processed'] += 1
            logger.info(f"Extracted {extraction_level} data from {file_info['name']}")
            
            return result
            
        except Exception as e:
            self.extraction_stats['errors'] += 1
            logger.error(f"Failed to extract data from {file_info.get('name', 'unknown')}: {str(e)}")
            return {
                'file_info': file_info,
                'extraction_level': extraction_level,
                'extraction_timestamp': datetime.utcnow().isoformat(),
                'extracted_data': None,
                'chunks': [],
                'metadata': {},
                'error': str(e)
            }
    
    def _determine_extraction_level(self, file_type: str, file_size: int) -> str:
        """
        Automatically determine the appropriate extraction level.
        
        Args:
            file_type: Type of file
            file_size: Size of file in bytes
            
        Returns:
            Extraction level to use
        """
        # Database files
        if file_type in ['sqlserver_db', 'oracle_db', 'sqlite_db', 'sql_dump']:
            return 'database'
        
        # Content files (if not too large)
        elif file_type in ['document', 'spreadsheet', 'presentation', 'log', 'config']:
            return 'content'
        
        # Large files or unknown types - metadata only
        elif file_size > self.max_file_size or file_type == 'unknown':
            return 'metadata'
        
        # Default to content for other types
        else:
            return 'content'
    
    def _extract_metadata_only(self, file_info: Dict[str, Any]) -> Dict[str, Any]:
        """
        Extract metadata only (Level 1).
        
        Args:
            file_info: File information
            
        Returns:
            Metadata dictionary
        """
        return {
            'level': 'metadata',
            'path': file_info.get('path', ''),
            'name': file_info.get('name', ''),
            'size': file_info.get('size', 0),
            'file_type': file_info.get('file_type', 'unknown'),
            'created_time': file_info.get('created_time'),
            'modified_time': file_info.get('modified_time'),
            'accessed_time': file_info.get('accessed_time'),
            'is_directory': file_info.get('is_directory', False),
            'extractable': file_info.get('extractable', False),
            'attributes': file_info.get('attributes', {})
        }
    
    def _extract_file_content(self, file_info: Dict[str, Any], unc_path: str) -> Dict[str, Any]:
        """
        Extract file content (Level 2).
        
        Args:
            file_info: File information
            unc_path: UNC path to the file
            
        Returns:
            Content extraction results
        """
        file_type = file_info.get('file_type', 'unknown')
        file_name = file_info.get('name', '')
        
        content_data = {
            'level': 'content',
            'file_type': file_type,
            'text_content': '',
            'metadata': {},
            'parsing_method': 'unknown'
        }
        
        try:
            # Read file content based on type
            if file_type == 'document':
                content_data.update(self._parse_document(unc_path))
            elif file_type == 'spreadsheet':
                content_data.update(self._parse_spreadsheet(unc_path))
            elif file_type == 'presentation':
                content_data.update(self._parse_presentation(unc_path))
            elif file_type in ['log', 'config', 'json', 'xml']:
                content_data.update(self._parse_text_file(unc_path))
            elif file_type == 'csv':
                content_data.update(self._parse_csv(unc_path))
            else:
                # Fallback to text parsing
                content_data.update(self._parse_text_file(unc_path))
                
        except Exception as e:
            logger.error(f"Failed to parse {file_name}: {str(e)}")
            content_data['error'] = str(e)
        
        return content_data
    
    def _parse_document(self, file_path: str) -> Dict[str, Any]:
        """Parse document files (PDF, DOCX, TXT)."""
        if not HAS_DOCUMENT_PARSERS:
            return {'text_content': '', 'parsing_method': 'unavailable', 'error': 'Document parsers not installed'}
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.pdf':
                return self._parse_pdf(file_path)
            elif file_ext == '.docx':
                return self._parse_docx(file_path)
            elif file_ext == '.txt':
                return self._parse_text_file(file_path)
            else:
                return {'text_content': '', 'parsing_method': 'unsupported', 'error': f'Unsupported document type: {file_ext}'}
                
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'error', 'error': str(e)}
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF files."""
        try:
            text_content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text_content += page.extract_text() + "\n"
            
            return {
                'text_content': text_content.strip(),
                'parsing_method': 'PyPDF2',
                'metadata': {'pages': len(pdf_reader.pages)}
            }
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'PyPDF2_error', 'error': str(e)}
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX files."""
        try:
            doc = Document(file_path)
            text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            return {
                'text_content': text_content.strip(),
                'parsing_method': 'python-docx',
                'metadata': {'paragraphs': len(doc.paragraphs)}
            }
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'python-docx_error', 'error': str(e)}
    
    def _parse_spreadsheet(self, file_path: str) -> Dict[str, Any]:
        """Parse spreadsheet files (XLSX, CSV)."""
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.xlsx':
                workbook = openpyxl.load_workbook(file_path)
                text_content = ""
                sheet_data = {}
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_text = ""
                    for row in sheet.iter_rows(values_only=True):
                        sheet_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
                    
                    text_content += f"Sheet: {sheet_name}\n{sheet_text}\n"
                    sheet_data[sheet_name] = {
                        'rows': sheet.max_row,
                        'columns': sheet.max_column
                    }
                
                return {
                    'text_content': text_content.strip(),
                    'parsing_method': 'openpyxl',
                    'metadata': {'sheets': sheet_data}
                }
            else:
                return self._parse_csv(file_path)
                
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'openpyxl_error', 'error': str(e)}
    
    def _parse_presentation(self, file_path: str) -> Dict[str, Any]:
        """Parse presentation files (PPTX)."""
        try:
            prs = Presentation(file_path)
            text_content = ""
            
            for i, slide in enumerate(prs.slides):
                text_content += f"Slide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
                text_content += "\n"
            
            return {
                'text_content': text_content.strip(),
                'parsing_method': 'python-pptx',
                'metadata': {'slides': len(prs.slides)}
            }
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'python-pptx_error', 'error': str(e)}
    
    def _parse_text_file(self, file_path: str) -> Dict[str, Any]:
        """Parse plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text_content = file.read()
            
            return {
                'text_content': text_content,
                'parsing_method': 'text_file'
            }
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'text_file_error', 'error': str(e)}
    
    def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        """Parse CSV files."""
        try:
            df = pd.read_csv(file_path)
            text_content = df.to_string(index=False)
            
            return {
                'text_content': text_content,
                'parsing_method': 'pandas_csv',
                'metadata': {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': list(df.columns)
                }
            }
        except Exception as e:
            return {'text_content': '', 'parsing_method': 'pandas_csv_error', 'error': str(e)}
    
    def _extract_database_content(self, file_info: Dict[str, Any], unc_path: str) -> Dict[str, Any]:
        """
        Extract database content (Level 3).
        
        Args:
            file_info: File information
            unc_path: UNC path to the database file
            
        Returns:
            Database extraction results
        """
        file_type = file_info.get('file_type', 'unknown')
        
        db_data = {
            'level': 'database',
            'file_type': file_type,
            'tables': [],
            'extraction_method': 'unknown',
            'error': None
        }
        
        try:
            if file_type == 'sqlite_db':
                db_data.update(self._extract_sqlite(unc_path))
            elif file_type == 'sql_dump':
                db_data.update(self._parse_sql_dump(unc_path))
            elif file_type in ['sqlserver_db', 'oracle_db']:
                db_data.update(self._extract_enterprise_db(file_info, unc_path))
            else:
                db_data['error'] = f"Unsupported database type: {file_type}"
                
        except Exception as e:
            logger.error(f"Failed to extract database content: {str(e)}")
            db_data['error'] = str(e)
        
        return db_data
    
    def _extract_sqlite(self, file_path: str) -> Dict[str, Any]:
        """Extract SQLite database content."""
        try:
            conn = sqlite3.connect(file_path)
            cursor = conn.cursor()
            
            # Get table names
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = [row[0] for row in cursor.fetchall()]
            
            table_data = {}
            for table_name in tables:
                # Get table schema
                cursor.execute(f"PRAGMA table_info({table_name});")
                columns = cursor.fetchall()
                
                # Get sample data (first 100 rows)
                cursor.execute(f"SELECT * FROM {table_name} LIMIT 100;")
                sample_data = cursor.fetchall()
                
                table_data[table_name] = {
                    'columns': [col[1] for col in columns],
                    'column_types': [col[2] for col in columns],
                    'sample_data': sample_data,
                    'total_rows': len(sample_data)
                }
            
            conn.close()
            
            return {
                'tables': table_data,
                'extraction_method': 'sqlite3',
                'table_count': len(tables)
            }
            
        except Exception as e:
            return {'tables': {}, 'extraction_method': 'sqlite3_error', 'error': str(e)}
    
    def _parse_sql_dump(self, file_path: str) -> Dict[str, Any]:
        """Parse SQL dump files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            # Extract table creation statements
            tables = {}
            lines = content.split('\n')
            
            current_table = None
            for line in lines:
                line = line.strip()
                if line.upper().startswith('CREATE TABLE'):
                    # Extract table name
                    table_name = line.split()[2].strip('`"\'')
                    current_table = table_name
                    tables[table_name] = {'schema': line, 'data': []}
                elif line.upper().startswith('INSERT INTO') and current_table:
                    tables[current_table]['data'].append(line)
            
            return {
                'tables': tables,
                'extraction_method': 'sql_dump_parser',
                'table_count': len(tables)
            }
            
        except Exception as e:
            return {'tables': {}, 'extraction_method': 'sql_dump_parser_error', 'error': str(e)}
    
    def _extract_enterprise_db(self, file_info: Dict[str, Any], unc_path: str) -> Dict[str, Any]:
        """Extract enterprise database content (SQL Server, Oracle)."""
        # This would require database-specific drivers and connection strings
        # For now, return a placeholder
        return {
            'tables': {},
            'extraction_method': 'enterprise_db_placeholder',
            'error': 'Enterprise database extraction requires additional configuration'
        }
    
    def _chunk_text(self, text: str) -> List[Dict[str, Any]]:
        """
        Chunk text into smaller pieces for ML processing.
        
        Args:
            text: Text to chunk
            
        Returns:
            List of text chunks with metadata
        """
        if not text:
            return []
        
        chunks = []
        words = text.split()
        
        for i in range(0, len(words), self.chunk_size):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = ' '.join(chunk_words)
            
            chunk = {
                'chunk_id': hashlib.md5(chunk_text.encode()).hexdigest()[:8],
                'text': chunk_text,
                'word_count': len(chunk_words),
                'char_count': len(chunk_text),
                'chunk_index': i // self.chunk_size,
                'start_word': i,
                'end_word': min(i + self.chunk_size, len(words))
            }
            chunks.append(chunk)
        
        self.extraction_stats['chunks_created'] += len(chunks)
        return chunks
    
    def get_extraction_stats(self) -> Dict[str, Any]:
        """Get extraction statistics."""
        return self.extraction_stats.copy()
    
    def reset_stats(self):
        """Reset extraction statistics."""
        self.extraction_stats = {
            'files_processed': 0,
            'chunks_created': 0,
            'databases_extracted': 0,
            'errors': 0
        }
